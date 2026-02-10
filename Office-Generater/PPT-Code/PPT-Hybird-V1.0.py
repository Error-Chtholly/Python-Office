import os
import pandas as pd
from pptx import Presentation
import re
from datetime import datetime
import warnings
import traceback
import copy
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk
import webbrowser
import urllib.parse
import ctypes
import math

warnings.filterwarnings('ignore')

# 尝试设置高DPI感知，修复模糊问题
try:
    ctypes.windll.shcore.SetProcessDpiAwareness(1)
except Exception:
    pass


class PPTGenerator:
    # ==========================================
    # 核心逻辑类 - 保持完全不变
    # ==========================================
    def __init__(self, template_path, excel_path, output_path, log_callback=None):
        self.template_path = template_path
        self.excel_path = excel_path
        self.output_path = output_path
        self.log_callback = log_callback
        self.template_pptx = None
        self.excel_data = None
        self.placeholders = set()

        self._load_template()
        self._load_excel_data()
        self._extract_placeholders()

    def log(self, message):
        print(message)
        if self.log_callback:
            self.log_callback(message)

    def _load_template(self):
        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"模板文件不存在: {self.template_path}")
        self.template_pptx = Presentation(self.template_path)
        self.log(f"成功加载模板: {self.template_path}")

    def _load_excel_data(self):
        if not os.path.exists(self.excel_path):
            raise FileNotFoundError(f"Excel文件不存在: {self.excel_path}")
        try:
            self.excel_data = pd.read_excel(self.excel_path, engine='openpyxl')
        except:
            self.excel_data = pd.read_excel(self.excel_path, engine='xlrd')

        self.excel_data.columns = self.excel_data.columns.str.strip()
        self.excel_data = self.excel_data.astype(str).apply(lambda x: x.str.strip())
        self.log(f"成功加载Excel数据，共 {len(self.excel_data)} 行")

    def _extract_placeholders(self):
        if len(self.template_pptx.slides) == 0:
            return
        slide = self.template_pptx.slides[0]
        pattern = r'\[([^\]]+)\]'
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                matches = re.findall(pattern, text)
                for match in matches:
                    self.placeholders.add(match.strip())
        self.log(f"检测到模板占位符: {list(self.placeholders)}")

    def _replace_text_in_shape(self, shape, replacements):
        if not hasattr(shape, "text_frame"):
            return False

        text_frame = shape.text_frame
        is_changed = False

        for paragraph in text_frame.paragraphs:
            original_text = paragraph.text
            if not any(f"[{k}]" in original_text for k in replacements):
                continue

            font_name = font_size = font_bold = font_italic = font_underline = font_color_rgb = None
            if len(paragraph.runs) > 0:
                ref_font = paragraph.runs[0].font
                font_name = ref_font.name
                font_size = ref_font.size
                font_bold = ref_font.bold
                font_italic = ref_font.italic
                font_underline = ref_font.underline
                try:
                    if hasattr(ref_font.color, 'rgb'):
                        font_color_rgb = ref_font.color.rgb
                except:
                    pass

            new_text = original_text
            sorted_keys = sorted(replacements.keys(), key=len, reverse=True)

            for placeholder in sorted_keys:
                value = replacements[placeholder]
                if value == "nan" or value is None:
                    value = ""
                pattern = r'\[' + re.escape(placeholder) + r'\]'
                new_text = re.sub(pattern, str(value), new_text)

            if new_text != original_text:
                paragraph.text = new_text
                if len(paragraph.runs) > 0:
                    new_run = paragraph.runs[0]
                    new_run.font.name = font_name
                    new_run.font.size = font_size
                    new_run.font.bold = font_bold
                    new_run.font.italic = font_italic
                    new_run.font.underline = font_underline
                    if font_color_rgb:
                        new_run.font.color.rgb = font_color_rgb
                is_changed = True
        return is_changed

    def run_single_mode(self):
        self.log("正在运行：Single 模式...")
        new_pptx = Presentation()
        new_pptx.slide_width = self.template_pptx.slide_width
        new_pptx.slide_height = self.template_pptx.slide_height

        slide_layout = self.template_pptx.slide_layouts[0]
        template_slide = self.template_pptx.slides[0]

        total = len(self.excel_data)
        for index, row in self.excel_data.iterrows():
            self.log(f"正在处理第 {index + 1}/{total} 行...")
            slide = new_pptx.slides.add_slide(slide_layout)

            for shape in list(slide.shapes):
                sp = shape._element
                sp.getparent().remove(sp)

            replacements = {}
            for col in self.excel_data.columns:
                val = row[col]
                replacements[col] = "" if val == "nan" else str(val)

            for shape in template_slide.shapes:
                try:
                    new_element = copy.deepcopy(shape._element)
                    slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
                except:
                    continue

            for shape in slide.shapes:
                try:
                    self._replace_text_in_shape(shape, replacements)
                except:
                    continue

        new_pptx.save(self.output_path)
        self.log(f"保存成功: {self.output_path}")

    def run_double_mode(self):
        self.log("正在运行：Double 模式...")
        new_pptx = Presentation()
        new_pptx.slide_width = self.template_pptx.slide_width
        new_pptx.slide_height = self.template_pptx.slide_height

        slide_layout = self.template_pptx.slide_layouts[0]
        template_slide = self.template_pptx.slides[0]
        columns = self.excel_data.columns
        total_rows = len(self.excel_data)

        for i in range(0, total_rows, 2):
            self.log(f"正在处理数据组: {i + 1} 和 {i + 2}...")
            slide = new_pptx.slides.add_slide(slide_layout)

            for shape in list(slide.shapes):
                sp = shape._element
                sp.getparent().remove(sp)

            replacements = {}
            row1 = self.excel_data.iloc[i]
            for col in columns:
                replacements[col] = row1[col]

            if i + 1 < total_rows:
                row2 = self.excel_data.iloc[i + 1]
                for col in columns:
                    replacements[f"{col}_2"] = row2[col]
            else:
                for col in columns:
                    replacements[f"{col}_2"] = ""

            for shape in template_slide.shapes:
                try:
                    new_element = copy.deepcopy(shape._element)
                    slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
                except:
                    continue

            for shape in slide.shapes:
                try:
                    self._replace_text_in_shape(shape, replacements)
                except:
                    continue

        new_pptx.save(self.output_path)
        self.log(f"保存成功: {self.output_path}")


class PPTToolGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("基于PPT和Excel的批量证书生成工具")

        # === 【修改】使用居中函数初始化窗口尺寸 ===
        # 原代码: self.root.geometry("900x1000")
        self._center_window(self.root, 900, 1000)

        # === 设置图标 ===
        try:
            if os.path.exists("logo.ico"):
                self.root.iconbitmap("logo.ico")
        except Exception:
            pass

        # === 二次元风格配置 ===
        self.font_main = ("Microsoft YaHei UI", 10)
        self.font_title = ("Microsoft YaHei UI", 11, "bold")
        self.font_bold = ("Microsoft YaHei UI", 10, "bold")
        self.font_link = ("Microsoft YaHei UI", 10, "underline")
        # 【新增】单选按钮字体
        self.font_radio = ("Microsoft YaHei UI", 11, "bold")

        self.bg_color = "#FFFBFD"
        self.accent_pink = "#FF85B3"
        self.accent_green = "#39C5BB"
        self.text_color = "#444444"

        self.root.configure(bg=self.bg_color)

        self._setup_styles()

        self.template_path = tk.StringVar()
        self.excel_path = tk.StringVar()
        self.output_path = tk.StringVar()
        self.is_double_mode = tk.BooleanVar(value=False)

        self._create_widgets()

    # === 【新增】窗口居中辅助函数 ===
    def _center_window(self, window, width, height):
        screen_width = window.winfo_screenwidth()
        screen_height = window.winfo_screenheight()
        x = int((screen_width - width) / 2)
        y = int((screen_height - height) / 2)
        window.geometry(f'{width}x{height}+{x}+{y}')

    def _setup_styles(self):
        style = ttk.Style()
        style.theme_use('clam')

        style.configure('TFrame', background=self.bg_color)
        style.configure('TLabel', background=self.bg_color, font=self.font_main, foreground=self.text_color)
        style.configure('TCheckbutton', background=self.bg_color, font=self.font_main, foreground=self.text_color)

        # 【新增】美化单选按钮 (Radiobutton) 样式
        style.configure('TRadiobutton',
                        background=self.bg_color,
                        font=self.font_radio,
                        foreground=self.text_color,
                        indicatorsize=16)  # 增大选择圈

        style.map('TRadiobutton',
                  foreground=[('active', self.accent_pink), ('selected', self.accent_pink)],
                  background=[('active', self.bg_color)],
                  indicatorcolor=[('selected', self.accent_pink), ('pressed', self.accent_pink)])

        style.configure('TLabelframe', background=self.bg_color, bordercolor=self.accent_pink)
        style.configure('TLabelframe.Label', background=self.bg_color, font=self.font_title,
                        foreground=self.accent_pink)

        style.configure('Accent.TButton', font=("Microsoft YaHei UI", 12, "bold"),
                        background=self.accent_green, foreground="white", borderwidth=0, padding=10)
        style.map('Accent.TButton',
                  background=[('active', self.accent_pink), ('pressed', '#FF69B4')],
                  foreground=[('active', 'white')])

        style.configure('Regular.TButton', font=self.font_main, background="#FFEFF5", foreground=self.accent_pink,
                        borderwidth=1, bordercolor=self.accent_pink)
        style.map('Regular.TButton',
                  background=[('active', self.accent_pink)],
                  foreground=[('active', 'white')])

        style.configure("Vertical.TScrollbar", background=self.bg_color, troughcolor="#FFF0F5",
                        bordercolor=self.bg_color, arrowcolor=self.accent_pink)

    def _create_widgets(self):
        main_frame = ttk.Frame(self.root, padding="25")
        main_frame.pack(fill='both', expand=True)

        # 1. 顶部标题
        title_lbl = tk.Label(main_frame, text="✨ 魔法证书生成工坊 ✨", font=("Microsoft YaHei UI", 18, "bold"),
                             bg=self.bg_color, fg=self.accent_green)
        title_lbl.pack(side='top', pady=(0, 20))

        # 2. 文件配置区域
        config_frame = ttk.LabelFrame(main_frame, text=" 📂 资源配置 (Files) ", padding="20")
        config_frame.pack(side='top', fill='x', pady=(0, 20))

        self._create_file_row(config_frame, "PPT 模板 (Template):", self.template_path)
        self._create_file_row(config_frame, "Excel 数据 (Data):", self.excel_path)
        self._create_file_row(config_frame, "保存位置 (Output):", self.output_path, is_save=True)

        # 3. 模式设置区域 (已美化)
        mode_frame = ttk.LabelFrame(main_frame, text=" ⚙️ 魔法阵列 (Layout Settings) ", padding="20")
        mode_frame.pack(side='top', fill='x', pady=(0, 20))

        # 提示文字
        tk.Label(mode_frame, text="请选择法阵的施法范围 (Target Scope)：",
                 font=self.font_bold, bg=self.bg_color, fg="#666").pack(anchor='w', pady=(0, 10))

        # 选项容器
        radio_frame = ttk.Frame(mode_frame)
        radio_frame.pack(fill='x', expand=True)

        # 选项1：单页单图
        rb1 = ttk.Radiobutton(
            radio_frame,
            text=" 📜 单页单图 (Single Mode) \n      [ 1 页PPT = 1 组数据 ]",
            variable=self.is_double_mode,
            value=False,
            cursor="hand2"
        )
        rb1.pack(side='left', padx=(0, 30))

        # 选项2：单页双图
        rb2 = ttk.Radiobutton(
            radio_frame,
            text=" 📑 单页双图 (Double Mode) \n      [ 1 页PPT = 2 组数据 ]",
            variable=self.is_double_mode,
            value=True,
            cursor="hand2"
        )
        rb2.pack(side='left')

        # 4. 运行按钮
        self.btn_run_text = tk.StringVar(value="✨ 启动魔法生成阵 (Start) ✨")
        self.btn_run = tk.Button(main_frame, textvariable=self.btn_run_text, command=self.run_generation,
                                 bg=self.accent_green, fg="white", font=("Microsoft YaHei UI", 14, "bold"),
                                 relief="flat", cursor="hand2", pady=10)
        self.btn_run.pack(side='top', fill='x', pady=(10, 20))

        self._animate_button()

        # 5. 底部信息区域
        bottom_frame = ttk.Frame(main_frame)
        bottom_frame.pack(side='bottom', fill='x', pady=(10, 5))

        self.status_label = tk.Label(bottom_frame, text="准备就绪... (Ready)", font=self.font_main, bg=self.bg_color,
                                     fg="#888")
        self.status_label.pack(side='top', pady=(0, 5))

        btn_container = ttk.Frame(bottom_frame)
        btn_container.pack(side='top', pady=(0, 8))

        btn_usage = ttk.Button(btn_container, text="使用说明 (Guide)", command=self.show_usage_info,
                               style='Regular.TButton', cursor="hand2")
        btn_usage.pack(side='left', padx=10)

        btn_about = ttk.Button(btn_container, text="关于软件 (About)", command=self.show_about_info,
                               style='Regular.TButton', cursor="hand2")
        btn_about.pack(side='left', padx=10)

        copyright_frame = ttk.Frame(bottom_frame)
        copyright_frame.pack(side='top', pady=5)

        ttk.Label(copyright_frame, text="© 所有版权归 Error Chtholly 所有（", foreground="#888").pack(side='left')
        link = ttk.Label(copyright_frame, text="https://github.com/error-chtholly",
                         foreground=self.accent_pink, cursor="hand2", font=self.font_link)
        link.pack(side='left')
        link.bind("<Button-1>", lambda e: self.open_link("https://github.com/error-chtholly"))
        ttk.Label(copyright_frame, text="）", foreground="#888").pack(side='left')

        # 6. 日志区域
        log_labelframe = ttk.LabelFrame(main_frame, text=" 📝 魔法咏唱日志 (Log) ", padding="15")
        log_labelframe.pack(side='top', fill='both', expand=True, pady=(0, 10))

        log_inner = ttk.Frame(log_labelframe)
        log_inner.pack(fill='both', expand=True)

        self.log_text = tk.Text(log_inner, height=12, state='disabled', bg="#FFFAFC", fg="#555",
                                font=("Consolas", 10), relief='flat', padx=10, pady=10)
        scrollbar = ttk.Scrollbar(log_inner, orient='vertical', command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=scrollbar.set)

        scrollbar.pack(side='right', fill='y')
        self.log_text.pack(side='left', fill='both', expand=True)

    def _animate_button(self):
        import time
        r_start, g_start, b_start = 57, 197, 187
        r_end, g_end, b_end = 255, 133, 179

        t = time.time() * 2
        factor = (math.sin(t) + 1) / 2

        r = int(r_start + (r_end - r_start) * factor)
        g = int(g_start + (g_end - g_start) * factor)
        b = int(b_start + (b_end - b_start) * factor)

        color_hex = f'#{r:02x}{g:02x}{b:02x}'

        if self.btn_run['state'] != 'disabled':
            self.btn_run.configure(bg=color_hex)

        self.root.after(50, self._animate_button)

    def _create_file_row(self, parent, label_text, var, is_save=False):
        row_frame = ttk.Frame(parent)
        row_frame.pack(fill='x', pady=8)

        ttk.Label(row_frame, text=label_text, width=18, anchor='e').pack(side='left', padx=(0, 10))

        entry = tk.Entry(row_frame, textvariable=var, font=self.font_main, bd=1, relief="solid", bg="white", fg="#555")
        entry.pack(side='left', fill='x', expand=True, padx=5, ipady=3)

        cmd = self.select_output if is_save else (self.select_template if "模板" in label_text else self.select_excel)
        ttk.Button(row_frame, text="📂 选择", command=cmd, style='Regular.TButton', width=8, cursor="hand2").pack(
            side='right')

    def open_link(self, url):
        webbrowser.open(url)

    def show_usage_info(self):
        top = tk.Toplevel(self.root)
        top.title("使用说明")
        # === 【修改】使用居中函数 ===
        self._center_window(top, 680, 700)
        top.resizable(False, False)
        top.configure(bg="#FFFBFD")

        # 头部
        header_frame = tk.Frame(top, bg=self.accent_pink, height=60)
        header_frame.pack(fill='x', side='top')
        tk.Label(header_frame, text="✨ 魔法吟唱指南 (Guide) ✨", font=("Microsoft YaHei UI", 16, "bold"),
                 bg=self.accent_pink, fg="white").pack(pady=15)

        # 内容
        content_frame = tk.Frame(top, bg="#FFFBFD", padx=30, pady=15)
        content_frame.pack(expand=True, fill='both')

        def add_step_title(text):
            tk.Label(content_frame, text=text, font=("Microsoft YaHei UI", 11, "bold"),
                     bg="#FFFBFD", fg=self.accent_green, anchor='w').pack(fill='x', pady=(12, 5))

        def add_step_desc(text):
            tk.Label(content_frame, text=text, font=("Microsoft YaHei UI", 10),
                     bg="#FFFBFD", fg="#555", justify="left", anchor='w').pack(fill='x')

        add_step_title("1. 准备魔力源 (Excel Data)")
        add_step_desc("创建一个 Excel 表格，第一行必须为【列名】（即动态变量名）。\n后续行为具体的动态数据。例如：包含“姓名”、“奖项”等列。")

        add_step_title("2. 绘制法阵 (PPT Template)")
        add_step_desc(
            "在 PPT 模板中，将需要动态变化的部分用方括号括起来。\n例如：输入 [姓名] 代表此处将自动替换为 Excel 中对应的姓名数据。\n固定不变的文字（如“荣誉证书”）直接保留即可。")

        add_step_title("3. 双重咏唱 (Double Mode - Optional)")
        add_step_desc(
            "如果需要在一页 PPT 上生成两个证书（左右或上下）：\n请在主界面选择“单页双图”模式。\n此时，第二组数据的占位符需要添加 _2 后缀（例如：[姓名_2]）。")

        add_step_title("4. 进阶秘籍 (Tutorial)")
        tk.Label(content_frame, text="更详细的图文教程，请查阅官方魔法书库：",
                 font=("Microsoft YaHei UI", 10), bg="#FFFBFD", fg="#555", anchor='w').pack(fill='x')

        link_repo = tk.Label(content_frame, text="👉 点击查看 GitHub 官方图文教程 👈",
                             font=("Microsoft YaHei UI", 10, "bold", "underline"),
                             fg=self.accent_pink, bg="#FFFBFD", cursor="hand2")
        link_repo.pack(pady=5, anchor='w')
        link_repo.bind("<Button-1>", lambda e: self.open_link("https://github.com/error-chtholly/Python-Office"))

        def on_enter(e): btn_close.config(bg="#FF69B4")

        def on_leave(e): btn_close.config(bg=self.accent_pink)

        btn_close = tk.Button(top, text="× 明白了 ( >ω< ) ×", command=top.destroy,
                              bg=self.accent_pink, fg="white",
                              activebackground="#FF69B4", activeforeground="white",
                              font=("Microsoft YaHei UI", 11, "bold"),
                              relief="flat", padx=35, pady=8, cursor="hand2")

        btn_close.bind("<Enter>", on_enter)
        btn_close.bind("<Leave>", on_leave)
        btn_close.pack(side='bottom', pady=(0, 25))

    def show_about_info(self):
        top = tk.Toplevel(self.root)
        top.title("关于软件")
        # === 【修改】使用居中函数 ===
        self._center_window(top, 680, 700)
        top.resizable(False, False)
        top.configure(bg="#FFFBFD")

        header_frame = tk.Frame(top, bg=self.accent_pink, height=60)
        header_frame.pack(fill='x', side='top')
        tk.Label(header_frame, text="✨ 批量证书生成工具 ✨", font=("Microsoft YaHei UI", 16, "bold"),
                 bg=self.accent_pink, fg="white").pack(pady=15)

        content_frame = tk.Frame(top, bg="#FFFBFD", padx=30, pady=15)
        content_frame.pack(expand=True, fill='both')

        def add_line(text, bold=False, color="#555", font_size=10):
            f = ("Microsoft YaHei UI", font_size, "bold" if bold else "normal")
            tk.Label(content_frame, text=text, font=f, bg="#FFFBFD", fg=color).pack(pady=2)

        add_line("版本：V1.0  |  构建：2026年2月10日", bold=True, color=self.accent_green, font_size=11)

        row_repo = tk.Frame(content_frame, bg="#FFFBFD")
        row_repo.pack(pady=(8, 5))
        tk.Label(row_repo, text="软件仓库地址：", font=("Microsoft YaHei UI", 10), bg="#FFFBFD").pack(side='left')
        link_repo = tk.Label(row_repo, text="https://github.com/error-chtholly/Python-Office",
                             font=self.font_link, fg=self.accent_pink, bg="#FFFBFD", cursor="hand2")
        link_repo.pack(side='left')
        link_repo.bind("<Button-1>", lambda e: self.open_link("https://github.com/error-chtholly/Python-Office"))

        tk.Label(content_frame, text="", bg="#FFFBFD", font=("Arial", 2)).pack()
        add_line("可通过官方存储库获取最新更新版本")
        add_line("所有版权归 Error Chtholly 所有", bold=True, font_size=11)
        add_line("云南师范大学地理学部、南京师范大学地理科学学院")

        desc_frame = tk.LabelFrame(content_frame, text=" 功能简介 ", bg="#FFFBFD", fg=self.accent_green,
                                   font=("Microsoft YaHei UI", 9, "bold"), bd=1, relief="solid")
        desc_frame.pack(fill='x', padx=10, pady=(15, 10))
        tk.Label(desc_frame, text="根据PPT模板的占位符结合Excel表格数据\n一键批量生成PPT证书，支持单页单图与双图模式。",
                 font=("Microsoft YaHei UI", 10), bg="#FFFBFD", fg="#666", pady=10).pack()

        tk.Label(content_frame, text="欢迎关注有关Python-Office处理工具合集得更多信息！",
                 font=("Microsoft YaHei UI", 10), bg="#FFFBFD").pack(pady=(5, 0))
        link_wel = tk.Label(content_frame, text="👉 点击直达官方主页 👈",
                            font=("Microsoft YaHei UI", 11, "bold", "underline"),
                            fg=self.accent_green, bg="#FFFBFD", cursor="hand2")
        link_wel.pack(pady=5)
        link_wel.bind("<Button-1>", lambda e: self.open_link("https://github.com/error-chtholly/Python-Office"))

        def on_enter(e): btn_close.config(bg="#FF69B4")

        def on_leave(e): btn_close.config(bg=self.accent_pink)

        btn_close = tk.Button(top, text="× 关 闭 ( >ω< ) ×", command=top.destroy,
                              bg=self.accent_pink, fg="white",
                              activebackground="#FF69B4", activeforeground="white",
                              font=("Microsoft YaHei UI", 11, "bold"),
                              relief="flat", padx=35, pady=8, cursor="hand2")

        btn_close.bind("<Enter>", on_enter)
        btn_close.bind("<Leave>", on_leave)
        btn_close.pack(side='bottom', pady=(0, 25))

    def append_log(self, message):
        self.log_text.config(state='normal')
        timestamp = datetime.now().strftime("[%H:%M:%S] ")
        self.log_text.insert('end', timestamp + str(message) + "\n")
        self.log_text.see('end')
        self.log_text.config(state='disabled')
        self.root.update()

    def select_template(self):
        filename = filedialog.askopenfilename(title="选择PPT模板", filetypes=[("PowerPoint", "*.pptx")])
        if filename: self.template_path.set(filename)

    def select_excel(self):
        filename = filedialog.askopenfilename(title="选择Excel文件", filetypes=[("Excel", "*.xlsx *.xls")])
        if filename: self.excel_path.set(filename)

    def select_output(self):
        default_name = f"Result_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filename = filedialog.asksaveasfilename(title="保存输出文件", initialfile=default_name,
                                                filetypes=[("PowerPoint", "*.pptx")])
        if filename: self.output_path.set(filename)

    def report_error(self, error_msg):
        is_send = messagebox.askyesno(
            "💔 哎呀，出错了",
            "程序运行过程中发生了意料之外的错误...\n\n是否将错误信息自动发送给官方进行反馈？\n(将调用系统默认邮件客户端)"
        )

        if is_send:
            recipient = "zhouzetong_rs@163.com"
            subject = "Python-Office Tool Error Report"
            body_text = f"User Feedback Error Report:\n\n{error_msg}"
            body = urllib.parse.quote(body_text)
            mailto_link = f"mailto:{recipient}?subject={subject}&body={body}"

            try:
                webbrowser.open(mailto_link)
                self.append_log(">>> 📧 已尝试调起邮件客户端发送反馈。")
            except Exception as e:
                messagebox.showerror("发送失败", f"无法调起邮件客户端，请手动发送至 {recipient}")

    def run_generation(self):
        t_path = self.template_path.get()
        e_path = self.excel_path.get()
        o_path = self.output_path.get()

        if not all([t_path, e_path, o_path]):
            messagebox.showwarning("提示", "⚠️ 请先完善所有文件路径！")
            return

        self.status_label.config(text="🔥 正在全力施法中... (Processing)", fg=self.accent_pink)
        self.btn_run.config(state='disabled', bg="#ccc")

        self.log_text.config(state='normal')
        self.log_text.delete(1.0, 'end')
        self.log_text.config(state='disabled')
        self.root.update()

        try:
            generator = PPTGenerator(t_path, e_path, o_path, log_callback=self.append_log)

            if self.is_double_mode.get():
                generator.run_double_mode()
                mode_text = "Double (双行)"
            else:
                generator.run_single_mode()
                mode_text = "Single (单行)"

            self.status_label.config(text="✨ 生成完成 (Success)", fg=self.accent_green)
            self.append_log(">>> ✨ 所有任务执行完毕 ✨ <<<")
            messagebox.showinfo("🎉 成功", f"PPT 生成成功！\n模式: {mode_text}\n路径: {o_path}")

        except Exception as e:
            self.status_label.config(text="💔 发生错误 (Error)", fg="red")
            error_trace = traceback.format_exc()
            self.append_log(f"运行出错: {str(e)}")
            self.append_log(error_trace)
            self.report_error(error_trace)

        finally:
            self.btn_run.config(state='normal')


def main():
    root = tk.Tk()
    app = PPTToolGUI(root)
    root.mainloop()


if __name__ == "__main__":
    main()