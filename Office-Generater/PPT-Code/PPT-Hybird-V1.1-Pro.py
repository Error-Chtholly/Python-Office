import os
import pandas as pd
from pptx import Presentation
import re
from datetime import datetime
import warnings
import traceback
import copy
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk
import webbrowser
import urllib.parse
import ctypes
import math

warnings.filterwarnings('ignore')

# 尝试设置高DPI感知，修复模糊问题
try:
    ctypes.windll.shcore.SetProcessDpiAwareness(1)
except Exception:
    pass


class PPTGenerator:
    # ==========================================
    # 核心逻辑类 (完全保持不变)
    # ==========================================
    def __init__(self, template_path, excel_path, output_path, log_callback=None):
        self.template_path = template_path
        self.excel_path = excel_path
        self.output_path = output_path
        self.log_callback = log_callback
        self.template_pptx = None
        self.excel_data = None
        self.placeholders = set()

        self._load_template()
        self._load_excel_data()
        self._extract_placeholders()

    def log(self, message):
        print(message)
        if self.log_callback:
            self.log_callback(message)

    def _load_template(self):
        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"模板文件不存在: {self.template_path}")
        self.template_pptx = Presentation(self.template_path)
        self.log(f"成功加载模板: {self.template_path}")

    def _load_excel_data(self):
        if not os.path.exists(self.excel_path):
            raise FileNotFoundError(f"Excel文件不存在: {self.excel_path}")
        try:
            self.excel_data = pd.read_excel(self.excel_path, engine='openpyxl')
        except:
            self.excel_data = pd.read_excel(self.excel_path, engine='xlrd')

        self.excel_data.columns = self.excel_data.columns.str.strip()
        self.excel_data = self.excel_data.astype(str).apply(lambda x: x.str.strip())
        self.log(f"成功加载Excel数据，共 {len(self.excel_data)} 行")

    def _extract_placeholders(self):
        if len(self.template_pptx.slides) == 0:
            return
        slide = self.template_pptx.slides[0]
        pattern = r'\[([^\]]+)\]'
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                matches = re.findall(pattern, text)
                for match in matches:
                    self.placeholders.add(match.strip())
        self.log(f"检测到模板占位符: {list(self.placeholders)}")

    def _replace_text_in_shape(self, shape, replacements):
        if not hasattr(shape, "text_frame"):
            return False

        text_frame = shape.text_frame
        is_changed = False

        for paragraph in text_frame.paragraphs:
            original_text = paragraph.text
            if not any(f"[{k}]" in original_text for k in replacements):
                continue

            # 保存格式
            font_name = font_size = font_bold = font_italic = font_underline = font_color_rgb = None
            if len(paragraph.runs) > 0:
                ref_font = paragraph.runs[0].font
                font_name = ref_font.name
                font_size = ref_font.size
                font_bold = ref_font.bold
                font_italic = ref_font.italic
                font_underline = ref_font.underline
                try:
                    if hasattr(ref_font.color, 'rgb'):
                        font_color_rgb = ref_font.color.rgb
                except:
                    pass

            new_text = original_text
            sorted_keys = sorted(replacements.keys(), key=len, reverse=True)

            for placeholder in sorted_keys:
                value = replacements[placeholder]
                if value == "nan" or value is None:
                    value = ""
                pattern = r'\[' + re.escape(placeholder) + r'\]'
                new_text = re.sub(pattern, str(value), new_text)

            if new_text != original_text:
                paragraph.text = new_text
                # 恢复格式
                if len(paragraph.runs) > 0:
                    new_run = paragraph.runs[0]
                    new_run.font.name = font_name
                    new_run.font.size = font_size
                    new_run.font.bold = font_bold
                    new_run.font.italic = font_italic
                    new_run.font.underline = font_underline
                    if font_color_rgb:
                        new_run.font.color.rgb = font_color_rgb
                is_changed = True
        return is_changed

    def run_general_mode(self, records_per_page=1):
        mode_name = "Single" if records_per_page == 1 else f"{records_per_page}-Up"
        self.log(f"正在运行：{mode_name} 融合模式 (每页 {records_per_page} 个)...")

        new_pptx = Presentation()
        new_pptx.slide_width = self.template_pptx.slide_width
        new_pptx.slide_height = self.template_pptx.slide_height

        slide_layout = self.template_pptx.slide_layouts[0]
        template_slide = self.template_pptx.slides[0]
        columns = self.excel_data.columns
        total_rows = len(self.excel_data)

        for i in range(0, total_rows, records_per_page):
            current_batch = (i // records_per_page) + 1
            total_batches = math.ceil(total_rows / records_per_page)
            self.log(
                f"正在处理页面: {current_batch}/{total_batches} (数据行 {i + 1}-{min(i + records_per_page, total_rows)})...")

            slide = new_pptx.slides.add_slide(slide_layout)

            for shape in list(slide.shapes):
                sp = shape._element
                sp.getparent().remove(sp)

            replacements = {}

            for offset in range(records_per_page):
                data_index = i + offset
                suffix = f"_{offset + 1}" if offset > 0 else ""

                if data_index < total_rows:
                    row = self.excel_data.iloc[data_index]
                    for col in columns:
                        val = row[col]
                        replacements[f"{col}{suffix}"] = "" if val == "nan" else str(val)
                else:
                    for col in columns:
                        replacements[f"{col}{suffix}"] = ""

            for shape in template_slide.shapes:
                try:
                    new_element = copy.deepcopy(shape._element)
                    slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
                except:
                    continue

            for shape in slide.shapes:
                try:
                    self._replace_text_in_shape(shape, replacements)
                except:
                    continue

        new_pptx.save(self.output_path)
        self.log(f"保存成功: {self.output_path}")


class PPTToolGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("基于PPT和Excel的批量证书生成工具 (Pro)")

        # === 【修改】使用居中函数初始化窗口尺寸 ===
        # 原代码: self.root.geometry("900x1000")
        self._center_window(self.root, 900, 1000)

        # === 设置图标 ===
        self._set_window_icon(self.root)

        # === 二次元风格配置 ===
        self.font_main = ("Microsoft YaHei UI", 10)
        self.font_title = ("Microsoft YaHei UI", 11, "bold")
        self.font_bold = ("Microsoft YaHei UI", 10, "bold")
        self.font_link = ("Microsoft YaHei UI", 10, "underline")

        # 【修改】修复：调小了复选框文字字号，使其看起来不那么离谱
        self.font_radio = ("Microsoft YaHei UI", 11, "bold")

        self.bg_color = "#FFFBFD"
        self.accent_pink = "#FF85B3"
        self.accent_green = "#39C5BB"
        self.text_color = "#444444"

        self.root.configure(bg=self.bg_color)

        self._setup_styles()

        self.template_path = tk.StringVar()
        self.excel_path = tk.StringVar()
        self.output_path = tk.StringVar()

        # === 模式选择变量 ===
        self.mode_var = tk.IntVar(value=1)
        self.custom_n_var = tk.StringVar(value="")

        self._create_widgets()

    # === 【新增】窗口居中辅助函数 ===
    def _center_window(self, window, width, height):
        # 获取屏幕宽度和高度
        screen_width = window.winfo_screenwidth()
        screen_height = window.winfo_screenheight()

        # 计算居中的 x 和 y 坐标
        x = int((screen_width - width) / 2)
        y = int((screen_height - height) / 2)

        # 设置窗口几何形状
        window.geometry(f'{width}x{height}+{x}+{y}')

    def _set_window_icon(self, window):
        try:
            if os.path.exists("logo.ico"):
                window.iconbitmap("logo.ico")
        except Exception:
            pass

    def _setup_styles(self):
        style = ttk.Style()
        style.theme_use('clam')

        style.configure('TFrame', background=self.bg_color)
        style.configure('TLabel', background=self.bg_color, font=self.font_main, foreground=self.text_color)

        # 【修改】修复：复选框样式调整
        # indicatorsize 调整为 16 (原20)，字体使用调整后的 font_radio
        style.configure('TRadiobutton',
                        background=self.bg_color,
                        font=self.font_radio,
                        foreground=self.text_color,
                        indicatorsize=16)

        style.map('TRadiobutton',
                  foreground=[('active', self.accent_pink), ('selected', self.accent_pink)],
                  background=[('active', self.bg_color)],
                  indicatorcolor=[('selected', self.accent_pink), ('pressed', self.accent_pink)])

        style.configure('TLabelframe', background=self.bg_color, bordercolor=self.accent_pink)
        style.configure('TLabelframe.Label', background=self.bg_color, font=self.font_title,
                        foreground=self.accent_pink)

        style.configure('Accent.TButton', font=("Microsoft YaHei UI", 12, "bold"),
                        background=self.accent_green, foreground="white", borderwidth=0, padding=10)
        style.map('Accent.TButton',
                  background=[('active', self.accent_pink), ('pressed', '#FF69B4')],
                  foreground=[('active', 'white')])

        style.configure('Regular.TButton', font=self.font_main, background="#FFEFF5", foreground=self.accent_pink,
                        borderwidth=1, bordercolor=self.accent_pink)
        style.map('Regular.TButton',
                  background=[('active', self.accent_pink)],
                  foreground=[('active', 'white')])

        style.configure("Vertical.TScrollbar", background=self.bg_color, troughcolor="#FFF0F5",
                        bordercolor=self.bg_color, arrowcolor=self.accent_pink)

    def _create_widgets(self):
        main_frame = ttk.Frame(self.root, padding="25")
        main_frame.pack(fill='both', expand=True)

        # 1. 顶部标题
        title_lbl = tk.Label(main_frame, text="✨ 魔法证书生成工坊 ✨", font=("Microsoft YaHei UI", 18, "bold"),
                             bg=self.bg_color, fg=self.accent_green)
        title_lbl.pack(side='top', pady=(0, 20))

        # 2. 文件配置区域
        config_frame = ttk.LabelFrame(main_frame, text=" 📂 资源配置 (Files) ", padding="20")
        config_frame.pack(side='top', fill='x', pady=(0, 20))

        self._create_file_row(config_frame, "PPT 模板 (Template):", self.template_path)
        self._create_file_row(config_frame, "Excel 数据 (Data):", self.excel_path)
        self._create_file_row(config_frame, "保存位置 (Output):", self.output_path, is_save=True)

        # 3. 模式设置区域
        mode_frame = ttk.LabelFrame(main_frame, text=" ⚙️ 魔法阵列 (Layout Settings) ", padding="20")
        mode_frame.pack(side='top', fill='x', pady=(0, 20))

        tk.Label(mode_frame, text="请选择一页PPT生成几个证书：", font=self.font_bold, bg=self.bg_color, fg="#666").pack(
            anchor='w', pady=(0, 10))

        radio_frame = ttk.Frame(mode_frame)
        radio_frame.pack(fill='x', expand=True)

        modes = [1, 2, 3, 4]
        for m in modes:
            rb = ttk.Radiobutton(radio_frame, text=f" {m} 个/页 ", variable=self.mode_var, value=m,
                                 command=self._on_mode_change, cursor="hand2")
            rb.pack(side='left', padx=(0, 20))

        rb_custom = ttk.Radiobutton(radio_frame, text=" 其他: ", variable=self.mode_var, value=-1,
                                    command=self._on_mode_change, cursor="hand2")
        rb_custom.pack(side='left', padx=(0, 5))

        # 自定义数量输入框
        custom_input_border = tk.Frame(radio_frame, bg=self.accent_pink, bd=0, padx=2, pady=2)
        custom_input_border.pack(side='left')

        self.entry_custom = tk.Entry(custom_input_border, textvariable=self.custom_n_var, width=5,
                                     font=("Microsoft YaHei UI", 12),
                                     bd=0, relief="flat", justify="center")
        self.entry_custom.pack(fill='both', expand=True)

        tk.Label(radio_frame, text=" 个/页", bg=self.bg_color, font=self.font_main).pack(side='left', padx=(5, 0))

        self._on_mode_change()

        # 4. 运行按钮
        self.btn_run_text = tk.StringVar(value="✨ 启动魔法生成阵 (Start) ✨")
        self.btn_run = tk.Button(main_frame, textvariable=self.btn_run_text, command=self.run_generation,
                                 bg=self.accent_green, fg="white", font=("Microsoft YaHei UI", 14, "bold"),
                                 relief="flat", cursor="hand2", pady=10)
        self.btn_run.pack(side='top', fill='x', pady=(10, 20))

        self._animate_button()

        # 5. 底部信息区域
        bottom_frame = ttk.Frame(main_frame)
        bottom_frame.pack(side='bottom', fill='x', pady=(10, 5))

        self.status_label = tk.Label(bottom_frame, text="准备就绪... (Ready)", font=self.font_main, bg=self.bg_color,
                                     fg="#888")
        self.status_label.pack(side='top', pady=(0, 5))

        btn_container = ttk.Frame(bottom_frame)
        btn_container.pack(side='top', pady=(0, 8))

        # === 修改处：添加访问网页按钮 (左侧) ===
        btn_web = ttk.Button(btn_container, text="访问网页 (Web)",
                             command=lambda: self.open_link("https://error-chtholly.github.io/Office-Generater/PPT-Code/V1.1-Pro.html"),
                             style='Regular.TButton', cursor="hand2")
        btn_web.pack(side='left', padx=10)

        btn_usage = ttk.Button(btn_container, text="使用说明 (Guide)", command=self.show_usage_info,
                               style='Regular.TButton', cursor="hand2")
        btn_usage.pack(side='left', padx=10)

        btn_about = ttk.Button(btn_container, text="关于软件 (About)", command=self.show_about_info,
                               style='Regular.TButton', cursor="hand2")
        btn_about.pack(side='left', padx=10)

        # === 修改处：添加退出应用按钮 (右侧) ===
        btn_exit = ttk.Button(btn_container, text="退出应用 (Exit)", command=self.root.quit,
                              style='Regular.TButton', cursor="hand2")
        btn_exit.pack(side='left', padx=10)

        copyright_frame = ttk.Frame(bottom_frame)
        copyright_frame.pack(side='top', pady=5)

        ttk.Label(copyright_frame, text="© 所有版权归 Error Chtholly 所有（", foreground="#888").pack(side='left')
        link = ttk.Label(copyright_frame, text="https://github.com/error-chtholly",
                         foreground=self.accent_pink, cursor="hand2", font=self.font_link)
        link.pack(side='left')
        link.bind("<Button-1>", lambda e: self.open_link("https://github.com/error-chtholly"))
        ttk.Label(copyright_frame, text="）", foreground="#888").pack(side='left')

        # 6. 日志区域
        log_labelframe = ttk.LabelFrame(main_frame, text=" 📝 魔法咏唱日志 (Log) ", padding="15")
        log_labelframe.pack(side='top', fill='both', expand=True, pady=(0, 10))

        log_inner = ttk.Frame(log_labelframe)
        log_inner.pack(fill='both', expand=True)

        self.log_text = tk.Text(log_inner, height=12, state='disabled', bg="#FFFAFC", fg="#555",
                                font=("Consolas", 10), relief='flat', padx=10, pady=10)
        scrollbar = ttk.Scrollbar(log_inner, orient='vertical', command=self.log_text.yview)
        self.log_text.configure(yscrollcommand=scrollbar.set)

        scrollbar.pack(side='right', fill='y')
        self.log_text.pack(side='left', fill='both', expand=True)

    def _on_mode_change(self):
        if self.mode_var.get() == -1:
            self.entry_custom.config(state='normal', bg="white")
            self.entry_custom.focus()
        else:
            self.entry_custom.config(state='disabled', bg="#f0f0f0")
            self.custom_n_var.set("")

    def _animate_button(self):
        import time
        r_start, g_start, b_start = 57, 197, 187
        r_end, g_end, b_end = 255, 133, 179

        t = time.time() * 2
        factor = (math.sin(t) + 1) / 2

        r = int(r_start + (r_end - r_start) * factor)
        g = int(g_start + (g_end - g_start) * factor)
        b = int(b_start + (b_end - b_start) * factor)

        color_hex = f'#{r:02x}{g:02x}{b:02x}'

        if self.btn_run['state'] != 'disabled':
            self.btn_run.configure(bg=color_hex)

        self.root.after(50, self._animate_button)

    def _create_file_row(self, parent, label_text, var, is_save=False):
        row_frame = ttk.Frame(parent)
        row_frame.pack(fill='x', pady=8)
        ttk.Label(row_frame, text=label_text, width=18, anchor='e').pack(side='left', padx=(0, 10))

        entry_border = tk.Frame(row_frame, bg=self.accent_pink, bd=0, padx=2, pady=2)
        entry_border.pack(side='left', fill='x', expand=True, padx=5)

        entry = tk.Entry(entry_border, textvariable=var, font=("Microsoft YaHei UI", 11),
                         bd=0, relief="flat", bg="white", fg="#555")
        entry.pack(fill='both', expand=True, ipady=4)

        cmd = self.select_output if is_save else (self.select_template if "模板" in label_text else self.select_excel)
        ttk.Button(row_frame, text="📂 选择", command=cmd, style='Regular.TButton', width=8, cursor="hand2").pack(
            side='right')

    def open_link(self, url):
        webbrowser.open(url)

    # ==========================================
    # 修复 & 美化：使用说明界面
    # ==========================================
    def show_usage_info(self):
        top = tk.Toplevel(self.root)
        top.title("使用说明")

        # === 【修改】使用居中函数 ===
        # 原代码: top.geometry("720x900")
        self._center_window(top, 720, 900)

        top.resizable(True, True)

        top.configure(bg="#FFFBFD")

        # 修复图标
        self._set_window_icon(top)

        # 头部
        header_frame = tk.Frame(top, bg=self.accent_pink, height=60)
        header_frame.pack(fill='x', side='top')
        tk.Label(header_frame, text="✨ 魔法吟唱指南 (Guide) ✨", font=("Microsoft YaHei UI", 16, "bold"),
                 bg=self.accent_pink, fg="white").pack(pady=15)

        # 内容容器
        main_content = tk.Frame(top, bg="#FFFBFD", padx=30, pady=20)
        main_content.pack(expand=True, fill='both')

        # 辅助函数：创建卡片式步骤
        def create_step_card(parent, number, title_text, desc_text):
            # 卡片边框
            card = tk.LabelFrame(parent, bg="#FFFBFD", bd=1, relief="solid",
                                 fg=self.accent_pink, font=("Microsoft YaHei UI", 10, "bold"),
                                 text=f" Step {number} ", padx=15, pady=10)
            # 【修改】减少卡片垂直间距，防止挤出
            card.pack(fill='x', pady=(0, 10))

            # 标题
            tk.Label(card, text=title_text, font=("Microsoft YaHei UI", 11, "bold"),
                     bg="#FFFBFD", fg=self.accent_green, anchor='w').pack(fill='x', pady=(0, 5))

            # 描述
            tk.Label(card, text=desc_text, font=("Microsoft YaHei UI", 10),
                     bg="#FFFBFD", fg="#666", justify="left", anchor='w').pack(fill='x')

        # 步骤 1
        create_step_card(main_content, "01", "准备魔力源 (Excel Data)",
                         "创建一个 Excel 表格，第一行必须为【列名】（变量名）。\n例如：包含“姓名”、“奖项”等列，后续行为具体数据。")

        # 步骤 2
        create_step_card(main_content, "02", "绘制法阵 (PPT Template)",
                         "在 PPT 模板中，用 [列名] 作为占位符。\n例如：输入 [姓名] 代表此处替换为 Excel 对应的姓名。")

        # 步骤 3
        create_step_card(main_content, "03", "多重影分身 (Multi-Layout Mode)",
                         "占位符命名规则：\n"
                         "• 第 1 个位置：[姓名]\n"
                         "• 第 2 个位置：[姓名_2]\n"
                         "• 第 N 个位置：[姓名_N]\n"
                         "（若某页数据不足，多余的占位符将自动置空）")

        # 步骤 4 (教程链接)
        tutorial_card = tk.LabelFrame(main_content, bg="#F0FDFC", bd=1, relief="solid", fg=self.accent_green,
                                      text=" 进阶秘籍 (Tutorial) ", font=("Microsoft YaHei UI", 10, "bold"), padx=15,
                                      pady=10)
        tutorial_card.pack(fill='x', pady=(0, 5))

        tk.Label(tutorial_card, text="更详细的图文教程，请查阅官方魔法书库：",
                 font=("Microsoft YaHei UI", 10), bg="#F0FDFC", fg="#555", anchor='w').pack(fill='x')

        link_repo = tk.Label(tutorial_card, text="👉 点击查看 GitHub 官方图文教程 👈",
                             font=("Microsoft YaHei UI", 11, "bold", "underline"),
                             fg=self.accent_pink, bg="#F0FDFC", cursor="hand2")
        link_repo.pack(pady=8, anchor='w')
        link_repo.bind("<Button-1>", lambda e: self.open_link("https://github.com/error-chtholly/Python-Office"))

        # 底部关闭按钮
        def on_enter(e): btn_close.config(bg="#FF69B4")

        def on_leave(e): btn_close.config(bg=self.accent_pink)

        btn_close = tk.Button(top, text="× 明白了 ( >ω< ) ×", command=top.destroy,
                              bg=self.accent_pink, fg="white",
                              activebackground="#FF69B4", activeforeground="white",
                              font=("Microsoft YaHei UI", 11, "bold"),
                              relief="flat", padx=35, pady=8, cursor="hand2")

        btn_close.bind("<Enter>", on_enter)
        btn_close.bind("<Leave>", on_leave)
        btn_close.pack(side='bottom', pady=(0, 25))

    def show_about_info(self):
        top = tk.Toplevel(self.root)
        top.title("关于软件")

        # === 【修改】使用居中函数 ===
        # 原代码: top.geometry("680x700")
        self._center_window(top, 680, 700)

        top.resizable(False, False)
        top.configure(bg="#FFFBFD")

        # 修复图标
        self._set_window_icon(top)

        header_frame = tk.Frame(top, bg=self.accent_pink, height=60)
        header_frame.pack(fill='x', side='top')
        tk.Label(header_frame, text="✨ 批量证书生成工具 ✨", font=("Microsoft YaHei UI", 16, "bold"),
                 bg=self.accent_pink, fg="white").pack(pady=15)

        content_frame = tk.Frame(top, bg="#FFFBFD", padx=30, pady=15)
        content_frame.pack(expand=True, fill='both')

        def add_line(text, bold=False, color="#555", font_size=10):
            f = ("Microsoft YaHei UI", font_size, "bold" if bold else "normal")
            tk.Label(content_frame, text=text, font=f, bg="#FFFBFD", fg=color).pack(pady=2)

        add_line("版本：V1.1 Pro  |  构建：2026年2月10日", bold=True, color=self.accent_green, font_size=11)

        row_repo = tk.Frame(content_frame, bg="#FFFBFD")
        row_repo.pack(pady=(8, 5))
        tk.Label(row_repo, text="软件仓库地址：", font=("Microsoft YaHei UI", 10), bg="#FFFBFD").pack(side='left')
        link_repo = tk.Label(row_repo, text="https://github.com/error-chtholly/Python-Office",
                             font=self.font_link, fg=self.accent_pink, bg="#FFFBFD", cursor="hand2")
        link_repo.pack(side='left')
        link_repo.bind("<Button-1>", lambda e: self.open_link("https://github.com/error-chtholly/Python-Office"))

        tk.Label(content_frame, text="", bg="#FFFBFD", font=("Arial", 2)).pack()
        add_line("可通过官方存储库获取最新更新版本")
        add_line("所有版权归 Error Chtholly 所有", bold=True, font_size=11)
        add_line("云南师范大学地理学部、南京师范大学地理科学学院")

        desc_frame = tk.LabelFrame(content_frame, text=" 功能简介 ", bg="#FFFBFD", fg=self.accent_green,
                                   font=("Microsoft YaHei UI", 9, "bold"), bd=1, relief="solid")
        desc_frame.pack(fill='x', padx=10, pady=(15, 10))

        tk.Label(desc_frame, text="根据PPT模板占位符结合Excel表格数据一键批量生成证书。\n\n"
                                  "✨ 特性升级：\n"
                                  "支持任意数量排版（N图/页）模式！\n"
                                  "只需在模板中设置 [Tag], [Tag_2]...[Tag_N] 即可。",
                 font=("Microsoft YaHei UI", 10), bg="#FFFBFD", fg="#666", pady=10, justify="left").pack()

        tk.Label(content_frame, text="欢迎关注有关Python-Office处理工具合集得更多信息！",
                 font=("Microsoft YaHei UI", 10), bg="#FFFBFD").pack(pady=(5, 0))
        link_wel = tk.Label(content_frame, text="👉 点击直达官方主页 👈",
                            font=("Microsoft YaHei UI", 11, "bold", "underline"),
                            fg=self.accent_green, bg="#FFFBFD", cursor="hand2")
        link_wel.pack(pady=5)
        link_wel.bind("<Button-1>", lambda e: self.open_link("https://github.com/error-chtholly/Python-Office"))

        def on_enter(e): btn_close.config(bg="#FF69B4")

        def on_leave(e): btn_close.config(bg=self.accent_pink)

        btn_close = tk.Button(top, text="× 关 闭 ( >ω< ) ×", command=top.destroy,
                              bg=self.accent_pink, fg="white",
                              activebackground="#FF69B4", activeforeground="white",
                              font=("Microsoft YaHei UI", 11, "bold"),
                              relief="flat", padx=35, pady=8, cursor="hand2")

        btn_close.bind("<Enter>", on_enter)
        btn_close.bind("<Leave>", on_leave)
        btn_close.pack(side='bottom', pady=(0, 25))

    def append_log(self, message):
        self.log_text.config(state='normal')
        timestamp = datetime.now().strftime("[%H:%M:%S] ")
        self.log_text.insert('end', timestamp + str(message) + "\n")
        self.log_text.see('end')
        self.log_text.config(state='disabled')
        self.root.update()

    def select_template(self):
        filename = filedialog.askopenfilename(title="选择PPT模板", filetypes=[("PowerPoint", "*.pptx *.ppt")])
        if filename: self.template_path.set(filename)

    def select_excel(self):
        filename = filedialog.askopenfilename(title="选择Excel文件", filetypes=[("Excel", "*.xlsx *.xls *.csv")])
        if filename: self.excel_path.set(filename)

    def select_output(self):
        default_name = f"Result_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filename = filedialog.asksaveasfilename(title="保存输出文件", initialfile=default_name,
                                                filetypes=[("PowerPoint", "*.pptx *.ppt")])
        if filename: self.output_path.set(filename)

    def report_error(self, error_msg):
        is_send = messagebox.askyesno(
            "💔 哎呀，出错了",
            "程序运行过程中发生了意料之外的错误...\n\n是否将错误信息自动发送给官方进行反馈？\n(将调用系统默认邮件客户端)"
        )

        if is_send:
            recipient = "zhouzetong_rs@163.com"
            subject = "Python-Office Tool Error Report"
            body_text = f"User Feedback Error Report:\n\n{error_msg}"
            body = urllib.parse.quote(body_text)
            mailto_link = f"mailto:{recipient}?subject={subject}&body={body}"

            try:
                webbrowser.open(mailto_link)
                self.append_log(">>> 📧 已尝试调起邮件客户端发送反馈。")
            except Exception as e:
                messagebox.showerror("发送失败", f"无法调起邮件客户端，请手动发送至 {recipient}")

    def run_generation(self):
        t_path = self.template_path.get()
        e_path = self.excel_path.get()
        o_path = self.output_path.get()

        if not all([t_path, e_path, o_path]):
            messagebox.showwarning("提示", "⚠️ 请先完善所有文件路径！")
            return

        # === 获取并验证 N ===
        mode_val = self.mode_var.get()
        records_per_page = 1

        if mode_val == -1:  # 自定义模式
            raw_n = self.custom_n_var.get().strip()
            if not raw_n.isdigit() or int(raw_n) <= 0:
                messagebox.showwarning("输入错误", "⚠️ 自定义数量必须是大于 0 的整数！")
                self.entry_custom.focus()
                return
            records_per_page = int(raw_n)
        else:
            records_per_page = mode_val

        self.status_label.config(text=f"🔥 正在施法 (N={records_per_page})... (Processing)", fg=self.accent_pink)
        self.btn_run.config(state='disabled', bg="#ccc")

        self.log_text.config(state='normal')
        self.log_text.delete(1.0, 'end')
        self.log_text.config(state='disabled')
        self.root.update()

        try:
            generator = PPTGenerator(t_path, e_path, o_path, log_callback=self.append_log)

            # 直接调用通用的生成函数
            generator.run_general_mode(records_per_page)

            self.status_label.config(text="✨ 生成完成 (Success)", fg=self.accent_green)
            self.append_log(">>> ✨ 所有任务执行完毕 ✨ <<<")
            messagebox.showinfo("🎉 成功", f"PPT 生成成功！\n模式: {records_per_page}个/页\n路径: {o_path}")

        except Exception as e:
            self.status_label.config(text="💔 发生错误 (Error)", fg="red")
            error_trace = traceback.format_exc()
            self.append_log(f"运行出错: {str(e)}")
            self.append_log(error_trace)
            self.report_error(error_trace)

        finally:
            self.btn_run.config(state='normal')


def main():
    root = tk.Tk()
    app = PPTToolGUI(root)
    root.mainloop()


if __name__ == "__main__":
    main()